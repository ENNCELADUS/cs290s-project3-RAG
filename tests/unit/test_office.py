from __future__ import annotations

from io import BytesIO
from pathlib import Path

from openpyxl import Workbook
from pptx import Presentation

from rag_collection.office import extract_office_text, office_environment_status


def test_extract_xlsx_text_from_workbook(tmp_path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Courses"
    sheet.append(["Course", "Instructor"])
    sheet.append(["CS101", "SIST Faculty"])
    buffer = BytesIO()
    workbook.save(buffer)

    result = extract_office_text(buffer.getvalue(), ".xlsx", tmp_path / "course.tmp")

    assert result.parser == "xlsx"
    assert "Courses" in result.text
    assert "CS101" in result.text
    assert "SIST Faculty" in result.text


def test_extract_pptx_text_from_presentation(tmp_path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "SIST Seminar"
    slide.placeholders[1].text = "Machine learning lecture"
    buffer = BytesIO()
    presentation.save(buffer)

    result = extract_office_text(buffer.getvalue(), ".pptx", tmp_path / "seminar.tmp")

    assert result.parser == "pptx"
    assert "SIST Seminar" in result.text
    assert "Machine learning lecture" in result.text


def test_office_environment_status_reports_installed_extractors() -> None:
    status = office_environment_status()

    assert status["has_openpyxl"]
    assert status["has_python_docx"]
    assert status["has_python_pptx"]
    assert status["has_xlrd"]
