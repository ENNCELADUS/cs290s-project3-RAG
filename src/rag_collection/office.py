from __future__ import annotations

import subprocess
import zipfile
from dataclasses import dataclass
from importlib.util import find_spec
from io import BytesIO
from pathlib import Path
from shutil import which
from typing import Any
from xml.etree import ElementTree

from openpyxl import load_workbook

from .chunking import normalize_text

OFFICE_EXTENSIONS = {".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx"}

OFFICE_CONTENT_TYPES = {
    "application/msword": ".doc",
    "application/vnd.ms-excel": ".xls",
    "application/vnd.ms-powerpoint": ".ppt",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
}


def office_environment_status() -> dict[str, bool]:
    return {
        "has_openpyxl": find_spec("openpyxl") is not None,
        "has_python_docx": find_spec("docx") is not None,
        "has_python_pptx": find_spec("pptx") is not None,
        "has_textutil": which("textutil") is not None,
        "has_xlrd": find_spec("xlrd") is not None,
    }


@dataclass(frozen=True)
class OfficeTextResult:
    text: str
    parser: str
    flags: list[str]


def extract_office_text(body: bytes, extension: str, temp_path: Path) -> OfficeTextResult:
    extension = extension.lower()
    if extension == ".docx":
        return OfficeTextResult(_extract_docx(body), "docx", [])
    if extension == ".xls":
        text = _extract_xls(body)
        if text:
            return OfficeTextResult(text, "xls", [])
        text = _extract_legacy_office_with_textutil(body, extension, temp_path)
        if text:
            return OfficeTextResult(text, "xls_textutil", [])
        return OfficeTextResult("", "xls_unsupported", ["unsupported_office_binary"])
    if extension == ".xlsx":
        return OfficeTextResult(_extract_xlsx(body), "xlsx", [])
    if extension == ".pptx":
        return OfficeTextResult(_extract_pptx(body), "pptx", [])
    if extension in {".doc", ".ppt"}:
        text = _extract_legacy_office_with_textutil(body, extension, temp_path)
        if text:
            return OfficeTextResult(text, f"{extension[1:]}_textutil", [])
        return OfficeTextResult("", f"{extension[1:]}_unsupported", ["unsupported_office_binary"])
    return OfficeTextResult("", "office_unsupported", ["unsupported_office_binary"])


def _extract_docx(body: bytes) -> str:
    try:
        from docx import Document

        document = Document(BytesIO(body))
        lines = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if values:
                    lines.append("\t".join(values))
        text = normalize_text("\n".join(lines))
        if text:
            return text
    except Exception:
        pass

    with zipfile.ZipFile(BytesIO(body)) as archive:
        text_parts = _xml_text_parts(archive.read("word/document.xml"))
    return normalize_text("\n".join(text_parts))


def _extract_xls(body: bytes) -> str:
    try:
        import xlrd

        workbook = xlrd.open_workbook(file_contents=body)
    except Exception:
        return ""

    lines: list[str] = []
    for sheet in workbook.sheets():
        lines.append(sheet.name)
        for row_index in range(sheet.nrows):
            values = [_cell_to_text(sheet.cell_value(row_index, col_index)) for col_index in range(sheet.ncols)]
            values = [value for value in values if value]
            if values:
                lines.append("\t".join(values))
    return normalize_text("\n".join(lines))


def _extract_xlsx(body: bytes) -> str:
    workbook = load_workbook(BytesIO(body), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in workbook.worksheets:
        lines.append(sheet.title)
        for row in sheet.iter_rows(values_only=True):
            values = [str(value) for value in row if value is not None and str(value).strip()]
            if values:
                lines.append("\t".join(values))
    workbook.close()
    return normalize_text("\n".join(lines))


def _extract_pptx(body: bytes) -> str:
    try:
        from pptx import Presentation

        presentation = Presentation(BytesIO(body))
        lines: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text)
        text = normalize_text("\n".join(lines))
        if text:
            return text
    except Exception:
        pass

    lines: list[str] = []
    with zipfile.ZipFile(BytesIO(body)) as archive:
        slide_names = sorted(
            name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )
        for slide_name in slide_names:
            lines.extend(_xml_text_parts(archive.read(slide_name)))
    return normalize_text("\n".join(lines))


def _extract_legacy_office_with_textutil(body: bytes, extension: str, temp_path: Path) -> str:
    input_path = temp_path.with_suffix(extension)
    output_path = temp_path.with_suffix(".txt")
    input_path.write_bytes(body)
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-output", str(output_path), str(input_path)],
            check=False,
            capture_output=True,
            timeout=20,
        )
        if result.returncode != 0 or not output_path.exists():
            return ""
        return normalize_text(output_path.read_text(encoding="utf-8", errors="replace"))
    except (OSError, subprocess.SubprocessError):
        return ""
    finally:
        input_path.unlink(missing_ok=True)
        output_path.unlink(missing_ok=True)


def _xml_text_parts(xml_bytes: bytes) -> list[str]:
    root = ElementTree.fromstring(xml_bytes)
    parts: list[str] = []
    for element in root.iter():
        if element.text and element.tag.rsplit("}", 1)[-1] in {"t", "instrText"}:
            parts.append(element.text)
    return parts


def _cell_to_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value).strip()
